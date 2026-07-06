"""
06_build_pptx.py
================
Genera la presentación ejecutiva (8 min) de la Tarea 3 en PowerPoint,
en español, a partir del informe final (Informe_3_UCBTDCM_Sonez_Arriagada.pdf)
y cumpliendo la pauta HW03 (pitch ejecutivo: hallazgos, escenarios,
recomendaciones; sin detalle de código).

Cifras de escenarios: versión MNL-atributos (Escenario A toro 34,3->59,7/67,0;
Escenario B kappa 0,52->1,23, +0,71 pp), que coincide con las FIGURAS del
informe (Fig. 3 y 4) y con la metodología declarada. Nota: las Tablas 6 y 7
del PDF quedaron con las cifras antiguas de M6 (50,5; +3,16) — inconsistencia
interna del informe; la presentación usa las de las figuras.

Salida: Presentacion_3_UCBTDCM_Sonez_Arriagada.pptx (raíz de Tarea 3)

Run:  python scripts/06_build_pptx.py
"""

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = Path(__file__).resolve().parent.parent
FIG = BASE / "outputs" / "figures"
LOGO = BASE / "outputs" / "Imagenes" / "logo_dii.png"
OUT = BASE / "Presentacion_3_UCBTDCM_Sonez_Arriagada.pptx"

# --- paleta "sushi premium" ---
INK    = RGBColor(0x23, 0x27, 0x2E)   # nori, fondo oscuro
CREAM  = RGBColor(0xF7, 0xF4, 0xEE)   # fondo claro contenido
CORAL  = RGBColor(0xE8, 0x6A, 0x5C)   # acento (salmón) — cifras clave
WASABI = RGBColor(0x6E, 0x9B, 0x54)   # acento secundario
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x23, 0x27, 0x2E)   # texto sobre claro
MUTED  = RGBColor(0x6B, 0x70, 0x78)   # texto secundario
LIGHT  = RGBColor(0xEC, 0xEA, 0xE4)   # texto sobre oscuro
MUTEDL = RGBColor(0xA9, 0xAD, 0xB3)   # texto secundario sobre oscuro

TITLE_FONT = "Cambria"
BODY_FONT  = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------
def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def _set(run, size, color, font=BODY_FONT, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    return tf


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def add_line(tf, text, size, color, font=BODY_FONT, bold=False, italic=False,
             first=False, space_after=6, space_before=0, align=PP_ALIGN.LEFT,
             line=1.05):
    p = para(tf, first)
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    r = p.add_run(); r.text = text
    _set(r, size, color, font, bold, italic)
    return p


def title(s, text, color=DARK, y=0.45, size=30, w=12.0):
    tf = textbox(s, Inches(0.7), Inches(y), Inches(w), Inches(1.0))
    add_line(tf, text, size, color, TITLE_FONT, bold=True, first=True, line=1.0)
    return tf


def dot(s, x, y, d=0.14, color=CORAL):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                           Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.fill.background(); c.shadow.inherit = False
    return c


def numchip(s, x, y, n, d=0.42, color=CORAL):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                           Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    add_line(tf, str(n), 18, WHITE, TITLE_FONT, bold=True, first=True,
             align=PP_ALIGN.CENTER)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return c


def stat(s, x, y, number, label, w=3.2, ncolor=CORAL, tcolor=MUTED,
         nsize=40, lsize=12.5):
    tf = textbox(s, Inches(x), Inches(y), Inches(w), Inches(1.3))
    add_line(tf, number, nsize, ncolor, TITLE_FONT, bold=True, first=True,
             space_after=1, line=1.0)
    add_line(tf, label, lsize, tcolor, BODY_FONT, space_after=0, line=1.0)
    return tf


def image_contain(s, path, x, y, w, h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nx = x + (w - nw) / 2
    ny = y + (h - nh) / 2
    return s.shapes.add_picture(path, Inches(nx), Inches(ny),
                                Inches(nw), Inches(nh))


def chip(s, x, y, w, h, color=WHITE, radius=0.08):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                           Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


# ===========================================================================
# SLIDE 1 — Portada (oscuro)
# ===========================================================================
s = slide(INK)
# chip blanco para el logo (logo es oscuro sobre transparente)
chip(s, 0.7, 0.6, 2.5, 0.85, WHITE)
image_contain(s, str(LOGO), 0.85, 0.72, 2.2, 0.6)

tf = textbox(s, Inches(0.7), Inches(2.35), Inches(11.9), Inches(2.2))
add_line(tf, "Preferencias de sushi en Japón", 44, WHITE, TITLE_FONT,
         bold=True, first=True, space_after=6, line=1.0)
add_line(tf, "Heterogeneidad del consumidor y estrategia para el segmento premium",
         20, CORAL, BODY_FONT, space_after=0, line=1.05)

tf = textbox(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.9))
add_line(tf, "Understanding Consumer Behavior through Discrete Choice Models  ·  HW3  ·  2026-1",
         13, MUTEDL, BODY_FONT, first=True, space_after=14)
add_line(tf, "Matías Arriagada R.      Vicente Soñez I.", 16, WHITE, BODY_FONT,
         bold=True, space_after=4)
add_line(tf, "Prof. Ph.D. Sebastián Astroza Tagle    —    6 de julio de 2026, Concepción, Chile",
         12, MUTEDL, BODY_FONT, space_after=0)

# ===========================================================================
# SLIDE 2 — Contexto y pregunta
# ===========================================================================
s = slide(CREAM)
title(s, "El desafío: un océano azul premium")
bullets = [
    ("En Chile el sushi se consume como comida rápida; falta una oferta de alta gama.", CORAL),
    ("Japón —mercado maduro— ofrece lecciones para construir ese segmento premium.", CORAL),
    ("¿Qué mueve la preferencia de sushi y a qué consumidor apuntar?", WASABI),
]
yb = 1.75
tf = textbox(s, Inches(1.05), Inches(1.7), Inches(5.75), Inches(3.4))
for i, (t, col) in enumerate(bullets):
    dot(s, 0.72, yb + i * 0.83 + 0.04, 0.13, col)
    add_line(tf, t, 16, DARK, BODY_FONT, space_after=18, line=1.12,
             first=(i == 0))
add_line(tf, "Datos: 5.000 consumidores · 10 tipos de sushi · ranking completo (Kamishima, 2003).",
         12.5, MUTED, BODY_FONT, italic=True, space_before=6, space_after=0)

stat(s, 1.05, 5.4, "34,3 %", "eligió toro (atún graso) como su #1",
     w=5.6, nsize=44)
image_contain(s, str(FIG / "choice_shares_setA.png"), 7.0, 1.7, 5.9, 4.9)

# ===========================================================================
# SLIDE 3 — Enfoque: 3 preguntas, 3 modelos
# ===========================================================================
s = slide(CREAM)
title(s, "Tres preguntas, tres modelos")
rows = [
    ("Logit Multinomial", "¿Quién se desvía del gusto promedio? Interacciones por edad y región."),
    ("Nested Logit", "¿Cómo sustituye el consumidor cuando su favorito no está?"),
    ("Clases Latentes", "¿Qué segmentos de consumidor —no observables— existen?"),
]
y = 1.95
for i, (h, d) in enumerate(rows):
    dot(s, 0.9, y + 0.06, 0.30, CORAL)
    tf = textbox(s, Inches(1.5), Inches(y - 0.12), Inches(11.0), Inches(1.2))
    add_line(tf, h, 22, DARK, TITLE_FONT, bold=True, first=True, space_after=2,
             line=1.0)
    add_line(tf, d, 16, MUTED, BODY_FONT, space_after=0, line=1.05)
    y += 1.45
tf = textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6))
add_line(tf, "Estimados en Python (scipy) y R/Apollo, con validación cruzada entre ambos motores.",
         12.5, MUTED, BODY_FONT, italic=True, first=True)

# ===========================================================================
# SLIDE 4 — Hallazgo 1: precio = señal de calidad (Veblen)
# ===========================================================================
s = slide(CREAM)
title(s, "El precio funciona como señal de calidad")
stat(s, 0.9, 1.75, "β precio = +0,58", "coeficiente positivo y significativo (z = 24)",
     w=5.8, nsize=34)
tf = textbox(s, Inches(0.9), Inches(3.1), Inches(5.9), Inches(3.2))
add_line(tf, "Contraintuitivo: a mayor precio, mayor preferencia. Los sushis más caros —toro y uni— son los más elegidos.",
         16, DARK, BODY_FONT, first=True, space_after=10, line=1.12)
add_line(tf, "Es el Efecto Veblen: en una encuesta sin restricción de presupuesto, el precio señala calidad, no costo.",
         16, DARK, BODY_FONT, space_after=12, line=1.12)
add_line(tf, "⇒  No descontar los productos estrella.", 16, WASABI, BODY_FONT,
         bold=True, space_after=0, line=1.1)
image_contain(s, str(FIG / "item_attribute_map.png"), 7.05, 1.65, 5.8, 5.0)

# ===========================================================================
# SLIDE 5 — Hallazgo 2: edad y grasa + Este/Oeste
# ===========================================================================
s = slide(CREAM)
title(s, "La edad valora la grasa; hay brecha Este–Oeste")
stat(s, 0.9, 1.75, "β grasa × edad = +0,089", "efecto positivo y robusto (z = 5,2)",
     w=6.0, nsize=30)
tf = textbox(s, Inches(0.9), Inches(3.05), Inches(5.9), Inches(3.4))
add_line(tf, "A mayor edad, más utilidad del sushi grasoso (toro, anago) —incluso controlando por el gusto medio de cada ítem.",
         16, DARK, BODY_FONT, first=True, space_after=12, line=1.12)
add_line(tf, "El Oeste prefiere sabores más ligeros: el efecto regional conjunto es significativo (p = 0,0004).",
         16, DARK, BODY_FONT, space_after=12, line=1.12)
add_line(tf, "Precio y grasa están muy correlacionados (0,82), lo que dificulta separar ambos canales.",
         14, MUTED, BODY_FONT, italic=True, space_after=0, line=1.1)
image_contain(s, str(FIG / "eastwest_preferences.png"), 7.05, 1.7, 5.8, 4.9)

# ===========================================================================
# SLIDE 6 — Hallazgo 3: Nested Logit, IIA
# ===========================================================================
s = slide(CREAM)
title(s, "Se rompe la IIA: sustitución dentro del nido")
stat(s, 0.9, 1.75, "λ akami = 0,32", "nido del atún: toro · maguro · tekka_maki (t = −26)",
     w=6.0, nsize=34)
tf = textbox(s, Inches(0.9), Inches(3.05), Inches(5.9), Inches(3.4))
add_line(tf, "Toro, maguro y tekka_maki son sustitutos ~3× más cercanos de lo que asume el MNL.",
         16, DARK, BODY_FONT, first=True, space_after=12, line=1.12)
add_line(tf, "Si falta el toro, el consumidor migra a otro atún —no a un marisco cualquiera.",
         16, DARK, BODY_FONT, space_after=12, line=1.12)
add_line(tf, "La independencia de alternativas irrelevantes se rechaza (LR = 573,5).",
         14, MUTED, BODY_FONT, italic=True, space_after=0, line=1.1)
image_contain(s, str(FIG / "substitution_dendrogram.png"), 7.05, 1.75, 5.8, 4.8)

# ===========================================================================
# SLIDE 7 — Hallazgo 4: tres segmentos latentes
# ===========================================================================
s = slide(CREAM)
title(s, "Tres segmentos latentes de consumidores")
# tabla
rows = [
    ("Segmento", "Share", "Firma de gusto", "Perfil"),
    ("Gourmet premium", "33,7 %", "uni, toro, ikura", "mayor edad · 51% mujeres"),
    ("Gusto ligero", "31,3 %", "ebi (rechaza uni)", "joven · Oeste · 64% mujeres"),
    ("Fanático del atún", "34,9 %", "toro, maguro, tekka", "hombre · Este"),
]
tbl_w = Inches(6.5)
table = s.shapes.add_table(4, 4, Inches(0.7), Inches(2.0), tbl_w,
                           Inches(3.6)).table
table.columns[0].width = Inches(2.1)
table.columns[1].width = Inches(1.05)
table.columns[2].width = Inches(1.85)
table.columns[3].width = Inches(1.5)
seg_colors = [None, CORAL, WASABI, RGBColor(0xC0, 0x39, 0x2B)]
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.margin_left = Pt(6); cell.margin_right = Pt(4)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = val
        if ri == 0:
            _set(r, 13, WHITE, BODY_FONT, bold=True)
            cell.fill.solid(); cell.fill.fore_color.rgb = INK
        else:
            _set(r, 13, DARK, BODY_FONT, bold=(ci == 0))
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xEF, 0xEC, 0xE5)
        if ci == 1 and ri > 0:
            _set(r, 14, seg_colors[ri], BODY_FONT, bold=True)
image_contain(s, str(FIG / "latent_class_3d.png"), 7.25, 1.55, 5.7, 5.2)
tf = textbox(s, Inches(0.7), Inches(5.9), Inches(6.5), Inches(0.9))
add_line(tf, "Cada consumidor coloreado por su clase más probable; los tres segmentos tienen tamaños casi iguales.",
         12, MUTED, BODY_FONT, italic=True, first=True, line=1.05)

# ===========================================================================
# SLIDE 8 — Escenario A: shock de precio del toro
# ===========================================================================
s = slide(CREAM)
title(s, "Escenario A — Shock de precio del toro (+40%)")
tf = textbox(s, Inches(0.7), Inches(1.65), Inches(5.7), Inches(3.0))
add_line(tf, "Por el efecto Veblen, encarecer el toro lo hace MÁS atractivo: su share sube de 34,3% a 67,0% (Nested Logit).",
         16, DARK, BODY_FONT, first=True, space_after=12, line=1.12)
add_line(tf, "El NL revela lo invisible al MNL: sus vecinos de nido se desploman mucho más —maguro 8,1→0,6% y tekka 2,3→0,2%.",
         16, DARK, BODY_FONT, space_after=0, line=1.12)
stat(s, 0.7, 5.15, "+33 pp", "share del toro (NL)", w=2.6, nsize=34)
stat(s, 3.5, 5.15, "−7,5 pp", "maguro, canibalizado (NL)", w=3.0, nsize=34,
     ncolor=RGBColor(0xC0, 0x39, 0x2B))
image_contain(s, str(FIG / "scenario_A_nl_mnl.png"), 6.7, 1.9, 6.2, 4.4)

# ===========================================================================
# SLIDE 9 — Escenario B: promoción saludable
# ===========================================================================
s = slide(CREAM)
title(s, "Escenario B — Promoción saludable (kappa-maki)")
tf = textbox(s, Inches(0.7), Inches(1.65), Inches(5.7), Inches(2.4))
add_line(tf, "Rebaja de −30% en precio y mayor disponibilidad (frecuencia 0,40 → 0,88).",
         16, DARK, BODY_FONT, first=True, space_after=10, line=1.12)
add_line(tf, "El motor del crecimiento es la DISPONIBILIDAD (β freq = +2,17), no el precio: bajar precio incluso resta (Veblen).",
         16, DARK, BODY_FONT, space_after=0, line=1.12)
stat(s, 0.7, 4.35, "×2,4", "share: 0,52% → 1,23%", w=2.7, nsize=40, ncolor=WASABI)
stat(s, 3.6, 4.35, "+0,9 pp", "jóvenes del Oeste (mayor respuesta)", w=3.1,
     nsize=32, ncolor=WASABI)
image_contain(s, str(FIG / "scenario_B_segments.png"), 6.7, 1.9, 6.2, 4.4)

# ===========================================================================
# SLIDE 10 — Recomendaciones
# ===========================================================================
s = slide(CREAM)
title(s, "Recomendaciones para el CEO")
recs = [
    ("Proteger márgenes", "Sin promociones agresivas en las estrellas (toro, uni): el precio vende calidad."),
    ("Blindar el nido del atún", "Ante quiebre de toro, asegurar maguro y tekka_maki para retener la demanda dentro del nido."),
    ("Priorizar al 'fanático del atún'", "34,9% del mercado y mayor ticket: abastecimiento premium y marketing exclusivo."),
    ("Saludables sin descuento", "Crecer por disponibilidad y visibilidad comercial, no por rebajas de precio."),
]
y = 1.85
for i, (h, d) in enumerate(recs, 1):
    numchip(s, 0.85, y + 0.02, i)
    tf = textbox(s, Inches(1.5), Inches(y - 0.1), Inches(11.1), Inches(1.1))
    add_line(tf, h, 19, DARK, TITLE_FONT, bold=True, first=True, space_after=2,
             line=1.0)
    add_line(tf, d, 15, MUTED, BODY_FONT, space_after=0, line=1.06)
    y += 1.25

# ===========================================================================
# SLIDE 11 — Cierre (oscuro)
# ===========================================================================
s = slide(INK)
tf = textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.6))
add_line(tf, "Del dato a la estrategia", 40, WHITE, TITLE_FONT, bold=True,
         first=True, space_after=14, line=1.0)
add_line(tf, "Precio como señal de calidad  ·  sustitución por nido  ·  segmentos definidos",
         19, CORAL, BODY_FONT, space_after=2, line=1.1)
add_line(tf, "→ una hoja de ruta premium accionable.", 19, LIGHT, BODY_FONT,
         space_after=0, line=1.1)
tf = textbox(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.5))
add_line(tf, "Gracias.", 22, WHITE, TITLE_FONT, bold=True, first=True,
         space_after=8)
add_line(tf, "Matías Arriagada R.  ·  Vicente Soñez I.      |      github.com/maty20arriagada/Tarea-3-Sushi",
         12.5, MUTEDL, BODY_FONT, space_after=0)

prs.save(str(OUT))
print(f"Guardado: {OUT.name}  ({len(prs.slides._sldIdLst)} slides)")
